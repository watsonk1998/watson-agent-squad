# Copyright (c) ModelScope Contributors. All rights reserved.
"""
Recursive directory pre-scanner with LLM-driven candidate discovery.

Performs a breadth-first scan of ``paths`` to collect comprehensive
file metadata (title, size, type, author, page-count, keywords, content
preview), and optionally reads full content of small files.  The LLM then
ranks the most promising document candidates for a given query.
This is a *zero-index* approach — no pre-built vector indices required.
"""
import asyncio
import json
import logging
import mimetypes
import os
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Tuple, Union

from sirchmunk.llm.openai_chat import OpenAIChat
from sirchmunk.utils.file_utils import fast_extract

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class FileCandidate:
    """Comprehensive metadata for a scanned file.

    Attributes:
        path: Absolute path.
        filename: Basename.
        extension: Lowercase file extension.
        size_bytes: File size in bytes.
        modified_at: Last modification time (ISO format).
        created_at: Creation time (ISO format).
        mime_type: Inferred MIME type.
        title: Document title (from PDF/DOCX metadata or first heading).
        author: Document author (from PDF/DOCX metadata).
        page_count: Number of pages (PDF) / slides (PPTX) / sheets (XLSX).
        encoding: Detected text encoding (text files only).
        line_count: Approximate line count (text files only).
        keywords: Keywords extracted from metadata or content.
        preview: First ~N chars of content for LLM triage.
        full_content: Complete file content (only for small files).
        content_loaded: Whether full_content was populated.
        relevance: LLM-assigned relevance (high / medium / low / None).
        reason: LLM-provided reason for relevance rating.
    """

    path: str
    filename: str
    extension: str = ""
    size_bytes: int = 0
    modified_at: Optional[str] = None
    created_at: Optional[str] = None
    mime_type: str = "application/octet-stream"
    title: str = ""
    author: str = ""
    page_count: int = 0
    encoding: str = ""
    line_count: int = 0
    keywords: List[str] = field(default_factory=list)
    preview: str = ""
    full_content: str = ""
    content_loaded: bool = False
    relevance: Optional[Literal["high", "medium", "low"]] = None
    reason: str = ""

    def to_dict(self) -> Dict[str, Any]:
        return {
            "path": self.path,
            "filename": self.filename,
            "extension": self.extension,
            "size_bytes": self.size_bytes,
            "modified_at": self.modified_at,
            "created_at": self.created_at,
            "mime_type": self.mime_type,
            "title": self.title,
            "author": self.author,
            "page_count": self.page_count,
            "encoding": self.encoding,
            "line_count": self.line_count,
            "keywords": self.keywords,
            "preview": self.preview[:300],
            "content_loaded": self.content_loaded,
            "relevance": self.relevance,
            "reason": self.reason,
        }

    def to_summary(self) -> str:
        """Compact text summary for LLM consumption."""
        parts = [f"- **{self.filename}** ({self.extension}, {self._human_size()})"]
        if self.title:
            parts.append(f"  Title: {self.title}")
        if self.author:
            parts.append(f"  Author: {self.author}")
        if self.page_count > 0:
            parts.append(f"  Pages/Sheets: {self.page_count}")
        if self.encoding:
            parts.append(f"  Encoding: {self.encoding}")
        if self.line_count > 0:
            parts.append(f"  Lines: ~{self.line_count}")
        if self.keywords:
            parts.append(f"  Keywords: {', '.join(self.keywords[:8])}")
        if self.modified_at:
            parts.append(f"  Modified: {self.modified_at[:10]}")
        if self.preview:
            # Show more preview for files with content loaded
            max_preview = 500 if self.content_loaded else 200
            clean = self.preview[:max_preview].replace(chr(10), ' ').strip()
            parts.append(f"  Preview: {clean}")
        return "\n".join(parts)

    def _human_size(self) -> str:
        """Format size_bytes as human-readable string."""
        if self.size_bytes < 1024:
            return f"{self.size_bytes}B"
        elif self.size_bytes < 1024 * 1024:
            return f"{self.size_bytes / 1024:.1f}KB"
        else:
            return f"{self.size_bytes / (1024 * 1024):.1f}MB"


@dataclass
class ScanResult:
    """Aggregated result of a directory scan.

    Attributes:
        candidates: All discovered file candidates.
        ranked_candidates: Candidates ranked by LLM (subset of candidates).
        total_files: Total number of files discovered.
        total_dirs: Total number of directories traversed.
        scan_duration_ms: Wall-clock time for the scan phase.
        rank_duration_ms: Wall-clock time for the LLM ranking phase.
    """

    candidates: List[FileCandidate] = field(default_factory=list)
    ranked_candidates: List[FileCandidate] = field(default_factory=list)
    total_files: int = 0
    total_dirs: int = 0
    scan_duration_ms: float = 0.0
    rank_duration_ms: float = 0.0

    @property
    def high_relevance(self) -> List[FileCandidate]:
        return [c for c in self.ranked_candidates if c.relevance == "high"]

    @property
    def medium_relevance(self) -> List[FileCandidate]:
        return [c for c in self.ranked_candidates if c.relevance == "medium"]


# ---------------------------------------------------------------------------
# Supported formats and MIME mapping
# ---------------------------------------------------------------------------

_SCANNABLE_EXTENSIONS = {
    # Text
    ".txt", ".md", ".rst", ".log", ".csv", ".tsv",
    # Code / config
    ".py", ".js", ".ts", ".json", ".yaml", ".yml", ".xml",
    ".toml", ".cfg", ".ini", ".conf", ".sh", ".bash",
    # Documents
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    # Web
    ".html", ".htm", ".css",
}

_TEXT_EXTENSIONS = {
    ".txt", ".md", ".rst", ".log", ".csv", ".py", ".json",
    ".yaml", ".yml", ".xml", ".html", ".htm", ".sh", ".toml",
    ".cfg", ".ini", ".conf", ".js", ".ts", ".css", ".bash",
    ".tsv",
}

_MIME_MAP = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".json": "application/json",
    ".xml": "application/xml",
    ".html": "text/html",
    ".htm": "text/html",
    ".csv": "text/csv",
}

# Files below this threshold are considered "small" and eligible for
# full content loading (default 100 KB).
_SMALL_FILE_THRESHOLD = 100 * 1024


# ---------------------------------------------------------------------------
# Scanner
# ---------------------------------------------------------------------------

class DirectoryScanner:
    """Recursively pre-scan directories and use LLM to discover candidates.

    Performs two phases:

    1. **Scan phase**: Fast filesystem walk + rich metadata extraction.
       For small files (< ``small_file_threshold`` bytes), full content
       is read into ``FileCandidate.full_content`` so the LLM can
       directly evaluate the file without a separate read step.
    2. **Rank phase**: LLM triage of candidates to identify most promising
       files.

    Args:
        llm: OpenAI-compatible chat client for LLM ranking.
        max_depth: Maximum recursion depth (default: 8).
        max_files: Maximum files to scan (default: 500).
        max_preview_chars: Characters to extract per file for preview (default: 800).
        small_file_threshold: Files smaller than this (bytes) get full content loaded (default: 100KB).
        max_workers: Thread pool size for parallel metadata extraction.
        exclude_patterns: Glob patterns to skip (e.g., ``["*.pyc", "__pycache__"]``).
    """

    DEFAULT_EXCLUDE = {
        "__pycache__", ".git", ".svn", "node_modules", ".idea",
        ".vscode", ".cache", ".tox", ".eggs", "*.egg-info",
        ".DS_Store", "Thumbs.db",
    }

    def __init__(
        self,
        llm: Optional[OpenAIChat] = None,
        max_depth: int = 8,
        max_files: int = 500,
        max_preview_chars: int = 800,
        small_file_threshold: int = _SMALL_FILE_THRESHOLD,
        max_workers: int = 8,
        exclude_patterns: Optional[List[str]] = None,
    ) -> None:
        self.llm = llm
        self.max_depth = max_depth
        self.max_files = max_files
        self.max_preview_chars = max_preview_chars
        self.small_file_threshold = small_file_threshold
        self.max_workers = max_workers
        self.exclude_patterns = set(exclude_patterns) if exclude_patterns else set()
        self.exclude_patterns.update(self.DEFAULT_EXCLUDE)

    # ---- Public API ----

    async def scan(
        self,
        paths: Union[str, Path, List[str], List[Path]],
    ) -> ScanResult:
        """Phase 1: Recursively walk directories and extract file metadata.

        For each file: filesystem stat, MIME type inference, and format-
        specific metadata (PDF pages/author, DOCX properties, etc.).
        Small files additionally have their full content loaded.

        Args:
            paths: One or more root directories to scan.

        Returns:
            ScanResult with populated ``candidates`` list.
        """
        if isinstance(paths, (str, Path)):
            paths = [paths]
        paths = [Path(p).resolve() for p in paths]

        t_start = datetime.now()
        result = ScanResult()

        # Collect all file paths
        all_files: List[Path] = []
        for root in paths:
            if root.is_file():
                all_files.append(root)
                continue
            if not root.is_dir():
                logger.warning(f"[DirScanner] Scan path not found: {root}")
                continue
            self._walk(root, all_files, result, depth=0)

        result.total_files = len(all_files)
        logger.info(f"[DirScanner] Found {result.total_files} files in {result.total_dirs} dirs")

        # Truncate if too many
        if len(all_files) > self.max_files:
            logger.warning(f"[DirScanner] Truncating to {self.max_files} files (from {len(all_files)})")
            all_files = all_files[: self.max_files]

        # Extract metadata in parallel
        loop = asyncio.get_event_loop()
        with ThreadPoolExecutor(max_workers=self.max_workers) as pool:
            candidates = await loop.run_in_executor(
                pool,
                lambda: self._extract_metadata_batch(all_files),
            )

        result.candidates = candidates

        content_loaded = sum(1 for c in candidates if c.content_loaded)
        result.scan_duration_ms = (datetime.now() - t_start).total_seconds() * 1000
        logger.info(
            f"[DirScanner] Scan complete: {len(candidates)} candidates "
            f"({content_loaded} fully loaded) in {result.scan_duration_ms:.0f}ms"
        )

        return result

    async def rank(
        self,
        query: str,
        scan_result: ScanResult,
        top_k: int = 20,
    ) -> ScanResult:
        """Phase 2: Use LLM to rank candidates by relevance to query.

        Args:
            query: User's search query.
            scan_result: Output from ``scan()``.
            top_k: Number of top candidates to include in LLM analysis.

        Returns:
            Updated ScanResult with ``ranked_candidates`` populated.
        """
        if not self.llm:
            logger.warning("[DirScanner] No LLM provided — skipping rank phase")
            scan_result.ranked_candidates = scan_result.candidates
            return scan_result

        t_start = datetime.now()

        # Build compact summary for LLM
        summaries = [c.to_summary() for c in scan_result.candidates[:top_k]]
        scan_text = "\n\n".join(summaries)

        prompt = self._build_rank_prompt(query, scan_text)

        try:
            response = await self.llm.achat(
                messages=[{"role": "user", "content": prompt}],
                stream=False,
            )
            ranked = self._parse_rank_response(response.content, scan_result.candidates)
            scan_result.ranked_candidates = ranked
        except Exception as exc:
            logger.error(f"[DirScanner] LLM ranking failed: {exc}")
            scan_result.ranked_candidates = scan_result.candidates

        scan_result.rank_duration_ms = (datetime.now() - t_start).total_seconds() * 1000
        logger.info(
            f"[DirScanner] Rank complete: {len(scan_result.ranked_candidates)} ranked "
            f"in {scan_result.rank_duration_ms:.0f}ms"
        )

        return scan_result

    async def scan_and_rank(
        self,
        query: str,
        paths: Union[str, Path, List[str], List[Path]],
        top_k: int = 20,
    ) -> ScanResult:
        """Convenience: run both scan and rank phases.

        Args:
            query: User's search query.
            paths: Directories to scan.
            top_k: Number of candidates for LLM ranking.

        Returns:
            ScanResult with both ``candidates`` and ``ranked_candidates``.
        """
        result = await self.scan(paths)
        return await self.rank(query, result, top_k=top_k)

    # ---- Filesystem walking ----

    def _walk(
        self,
        root: Path,
        out: List[Path],
        result: ScanResult,
        depth: int,
    ) -> None:
        """Recursive directory walk with depth limiting and exclusion."""
        if depth > self.max_depth:
            return
        if len(out) >= self.max_files:
            return

        result.total_dirs += 1

        try:
            entries = sorted(root.iterdir(), key=lambda p: p.name)
        except PermissionError:
            logger.debug(f"[DirScanner] Permission denied: {root}")
            return

        for entry in entries:
            if len(out) >= self.max_files:
                return

            name = entry.name

            # Skip excluded patterns
            if name in self.exclude_patterns:
                continue
            if any(entry.match(pat) for pat in self.exclude_patterns if "*" in pat):
                continue
            # Skip hidden files/dirs
            if name.startswith("."):
                continue

            if entry.is_dir():
                self._walk(entry, out, result, depth + 1)
            elif entry.is_file():
                ext = entry.suffix.lower()
                if ext in _SCANNABLE_EXTENSIONS:
                    out.append(entry)

    # ---- Metadata extraction ----

    def _extract_metadata_batch(self, files: List[Path]) -> List[FileCandidate]:
        """Extract metadata for a batch of files (runs in thread pool)."""
        candidates: List[FileCandidate] = []
        for f in files:
            try:
                candidate = self._extract_single(f)
                if candidate:
                    candidates.append(candidate)
            except Exception as exc:
                logger.debug(f"[DirScanner] Metadata error for {f}: {exc}")
        return candidates

    def _extract_single(self, file_path: Path) -> Optional[FileCandidate]:
        """Extract comprehensive metadata from a single file.

        Extracts:
        - Filesystem stat (size, mtime, ctime)
        - MIME type
        - Format-specific metadata (PDF: title/author/pages, DOCX: props, etc.)
        - Content preview (first N chars)
        - Full content for small files (< small_file_threshold)
        """
        stat = file_path.stat()
        ext = file_path.suffix.lower()

        candidate = FileCandidate(
            path=str(file_path),
            filename=file_path.name,
            extension=ext,
            size_bytes=stat.st_size,
            modified_at=datetime.fromtimestamp(stat.st_mtime).isoformat(),
            created_at=datetime.fromtimestamp(stat.st_ctime).isoformat(),
            mime_type=_MIME_MAP.get(
                ext,
                mimetypes.guess_type(str(file_path))[0] or "application/octet-stream"
            ),
        )

        # Extract format-specific metadata + preview
        try:
            if ext == ".pdf":
                self._extract_pdf_metadata(file_path, candidate)
            elif ext in (".docx", ".pptx", ".xlsx"):
                self._extract_office_metadata(file_path, candidate)
            elif ext in _TEXT_EXTENSIONS:
                self._extract_text_metadata(file_path, candidate)
        except Exception as exc:
            logger.debug(f"[DirScanner] Content extraction failed for {file_path}: {exc}")

        # Full content loading for small files the LLM can digest directly
        if stat.st_size <= self.small_file_threshold and not candidate.content_loaded:
            self._try_load_full_content(file_path, candidate)

        return candidate

    # ---- Text files ----

    def _extract_text_metadata(self, path: Path, candidate: FileCandidate) -> None:
        """Extract preview, title, encoding, and line count from plain-text files."""
        # Detect encoding
        try:
            raw_head = path.read_bytes()[:4096]
            candidate.encoding = self._detect_encoding(raw_head)
        except Exception:
            candidate.encoding = "utf-8"

        try:
            with open(path, "r", encoding=candidate.encoding or "utf-8", errors="replace") as f:
                head = f.read(self.max_preview_chars + 500)
        except Exception:
            return

        lines = head.split("\n")

        # Title: first non-empty line (strip markdown heading markers)
        for line in lines:
            stripped = line.strip().lstrip("#").strip()
            if stripped and len(stripped) > 2:
                candidate.title = stripped[:120]
                break

        # Preview: first N chars
        candidate.preview = head[: self.max_preview_chars].strip()

        # Approximate line count (fast: count newlines in head, extrapolate)
        if candidate.size_bytes > 0:
            head_lines = head.count("\n")
            head_len = len(head.encode("utf-8", errors="replace"))
            if head_len > 0:
                candidate.line_count = int(
                    head_lines * (candidate.size_bytes / head_len)
                )

        # Keywords: extract from first ~10 meaningful lines
        keywords: List[str] = []
        for line in lines[:10]:
            stripped = line.strip()
            if len(stripped) > 5:
                tokens = [t.strip(".,;:!?()[]{}\"'") for t in stripped.split()]
                for t in tokens:
                    if len(t) > 3 and t.lower() not in keywords and len(keywords) < 10:
                        keywords.append(t.lower())
        candidate.keywords = keywords

    # ---- PDF ----

    def _extract_pdf_metadata(self, path: Path, candidate: FileCandidate) -> None:
        """Extract title, author, page count, keywords from PDF metadata + first page text."""
        try:
            import pypdf
        except ImportError:
            return

        import logging as _logging
        import warnings

        # Suppress noisy pypdf warnings for malformed PDF cross-references
        # pypdf uses logging.getLogger("pypdf._reader").warning(...)
        _pypdf_logger = _logging.getLogger("pypdf._reader")
        _prev_level = _pypdf_logger.level
        _pypdf_logger.setLevel(_logging.ERROR)

        try:
            with warnings.catch_warnings():
                warnings.filterwarnings("ignore", message=".*wrong pointing object.*")
                warnings.filterwarnings("ignore", message=".*Ignoring wrong.*")

                with open(path, "rb") as f:
                    reader = pypdf.PdfReader(f)
                    meta = reader.metadata or {}

                    candidate.title = (meta.get("/Title") or "")[:120]
                    candidate.author = (meta.get("/Author") or "")[:80]
                    candidate.page_count = len(reader.pages)

                    raw_keywords = meta.get("/Keywords") or ""
                    if raw_keywords:
                        import re
                        candidate.keywords = [
                            k.strip() for k in re.split(r"[,;|]", raw_keywords) if k.strip()
                        ]
                    if candidate.author and candidate.author not in candidate.keywords:
                        candidate.keywords.insert(0, candidate.author)

                    # Preview: first page text (more generous for PDFs)
                    if len(reader.pages) > 0:
                        text = reader.pages[0].extract_text() or ""
                        candidate.preview = text[: self.max_preview_chars].strip()

                        # If the file is small enough, extract all pages
                        if candidate.size_bytes <= self.small_file_threshold:
                            all_text_parts = []
                            for page in reader.pages:
                                page_text = page.extract_text() or ""
                                if page_text.strip():
                                    all_text_parts.append(page_text)
                            if all_text_parts:
                                candidate.full_content = "\n\n".join(all_text_parts)
                                candidate.content_loaded = True

        except Exception as exc:
            logger.debug(f"[DirScanner] PDF metadata error: {exc}")
        finally:
            _pypdf_logger.setLevel(_prev_level)

    # ---- Office documents ----

    def _extract_office_metadata(self, path: Path, candidate: FileCandidate) -> None:
        """Extract metadata from Office documents (docx, pptx, xlsx)."""
        ext = path.suffix.lower()
        try:
            if ext == ".docx":
                self._extract_docx_metadata(path, candidate)
            elif ext == ".xlsx":
                self._extract_xlsx_metadata(path, candidate)
            elif ext == ".pptx":
                self._extract_pptx_metadata(path, candidate)
        except Exception as exc:
            logger.debug(f"[DirScanner] Office metadata error: {exc}")

    def _extract_docx_metadata(self, path: Path, candidate: FileCandidate) -> None:
        """Extract metadata and content from DOCX files."""
        try:
            from docx import Document
        except ImportError:
            return

        doc = Document(str(path))
        props = doc.core_properties
        candidate.title = (props.title or "")[:120]
        candidate.author = (props.author or "")[:80]

        if props.keywords:
            candidate.keywords.extend(
                k.strip() for k in props.keywords.split(",") if k.strip()
            )
        if candidate.author and candidate.author not in candidate.keywords:
            candidate.keywords.insert(0, candidate.author)

        # Extract paragraph text
        para_texts: List[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                para_texts.append(text)
        candidate.line_count = len(para_texts)

        all_text = "\n".join(para_texts)
        candidate.preview = all_text[: self.max_preview_chars]

        # Full content for small files
        if candidate.size_bytes <= self.small_file_threshold and all_text:
            candidate.full_content = all_text
            candidate.content_loaded = True

    def _extract_xlsx_metadata(self, path: Path, candidate: FileCandidate) -> None:
        """Extract sheet names and sample data from Excel files."""
        try:
            import openpyxl
        except ImportError:
            return

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        candidate.keywords = list(wb.sheetnames)[:10]
        candidate.page_count = len(wb.sheetnames)
        candidate.title = f"Excel: {', '.join(wb.sheetnames[:3])}"

        # Preview: headers + first few rows of the first sheet
        try:
            ws = wb[wb.sheetnames[0]]
            preview_parts: List[str] = []
            for i, row in enumerate(ws.iter_rows(max_row=5, values_only=True)):
                cells = [str(c) if c is not None else "" for c in row]
                preview_parts.append(" | ".join(cells))
            candidate.preview = "\n".join(preview_parts)[: self.max_preview_chars]
        except Exception:
            pass

        wb.close()

    def _extract_pptx_metadata(self, path: Path, candidate: FileCandidate) -> None:
        """Extract slide count, titles, and text from PowerPoint files."""
        try:
            from pptx import Presentation
        except ImportError:
            return

        prs = Presentation(str(path))
        candidate.page_count = len(prs.slides)
        candidate.title = f"Slides: {len(prs.slides)} slides"

        # Preview: text from first few slides
        preview_parts: List[str] = []
        for slide in prs.slides[:5]:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    preview_parts.append(shape.text.strip())
        all_text = "\n".join(preview_parts)
        candidate.preview = all_text[: self.max_preview_chars]

        if candidate.size_bytes <= self.small_file_threshold and all_text:
            candidate.full_content = all_text
            candidate.content_loaded = True

    # ---- Full content loading ----

    def _try_load_full_content(self, path: Path, candidate: FileCandidate) -> None:
        """Load full file content for small files the LLM can digest.

        Only loads if the file is a text-like format and within the
        size threshold.  Binary formats (PDF, Office) are handled in
        their respective extractors above.
        """
        if candidate.content_loaded:
            return

        ext = path.suffix.lower()
        if ext not in _TEXT_EXTENSIONS:
            return

        try:
            enc = candidate.encoding or "utf-8"
            content = path.read_text(encoding=enc, errors="replace")
            candidate.full_content = content
            candidate.content_loaded = True
        except Exception as exc:
            logger.debug(f"[DirScanner] Full content load failed for {path}: {exc}")

    # ---- Helpers ----

    @staticmethod
    def _detect_encoding(raw: bytes) -> str:
        """Detect text encoding from raw bytes."""
        try:
            import charset_normalizer
            result = charset_normalizer.from_bytes(raw).best()
            return result.encoding if result else "utf-8"
        except ImportError:
            return "utf-8"

    # ---- LLM ranking ----

    @staticmethod
    def _build_rank_prompt(query: str, scan_text: str) -> str:
        """Build the LLM prompt for candidate ranking."""
        return f"""You are a document triage specialist. Analyze the directory scan results below and identify the most relevant files for answering the user's query.

## User Query
{query}

## Directory Scan Results
{scan_text}

## Instructions
1. Rank all scanned files by their likely relevance to the query.
2. For each file, provide a brief reason why it may or may not be relevant.
3. Return ONLY a JSON array — no other text.

## Output Format
```json
[
  {{"path": "/abs/path/to/file", "relevance": "high", "reason": "brief reason"}},
  {{"path": "/abs/path/to/file", "relevance": "medium", "reason": "brief reason"}},
  {{"path": "/abs/path/to/file", "relevance": "low", "reason": "brief reason"}}
]
```"""

    def _parse_rank_response(
        self,
        llm_response: str,
        candidates: List[FileCandidate],
    ) -> List[FileCandidate]:
        """Parse LLM ranking response and update candidate relevance.

        The LLM may return paths in any form — absolute, relative,
        basename-only, or with partial directory components.  We match
        each returned path against the candidate list using a cascade:
        1. Exact match on ``FileCandidate.path`` (absolute).
        2. Match on filename (basename) only.
        3. Match on ``endswith`` (partial path suffix).
        """
        # Build multiple lookup indices for flexible matching
        path_exact: Dict[str, FileCandidate] = {c.path: c for c in candidates}
        name_map: Dict[str, List[FileCandidate]] = {}
        for c in candidates:
            name_map.setdefault(c.filename, []).append(c)

        def _resolve(p: str) -> Optional[FileCandidate]:
            """Resolve an LLM-returned path to a FileCandidate."""
            # 1) Exact absolute match
            if p in path_exact:
                return path_exact[p]

            # 2) Basename match (e.g. "report.md" → "/full/path/report.md")
            basename = os.path.basename(p)
            hits = name_map.get(basename)
            if hits and len(hits) == 1:
                return hits[0]

            # 3) Suffix match (e.g. "DINOv3_zh/report.md" → "/data/DINOv3_zh/report.md")
            p_normalized = p.replace("\\", "/")
            for c in candidates:
                if c.path.replace("\\", "/").endswith(p_normalized):
                    return c

            # 4) If multiple basename hits, pick first
            if hits:
                return hits[0]

            return None

        ranked: List[FileCandidate] = []
        matched_paths: set = set()
        try:
            import re
            json_match = re.search(r"\[.*\]", llm_response, re.DOTALL)
            if not json_match:
                return candidates

            items = json.loads(json_match.group())
            for item in items:
                p = item.get("path", "")
                relevance = item.get("relevance", "low")
                reason = item.get("reason", "")

                c = _resolve(p)
                if c and c.path not in matched_paths:
                    c.relevance = relevance if relevance in ("high", "medium", "low") else "low"
                    c.reason = reason
                    ranked.append(c)
                    matched_paths.add(c.path)

        except (json.JSONDecodeError, Exception) as exc:
            logger.warning(f"[DirScanner] Failed to parse rank response: {exc}")
            return candidates

        # Sort: high > medium > low
        relevance_order = {"high": 0, "medium": 1, "low": 2}
        ranked.sort(key=lambda c: relevance_order.get(c.relevance, 3))

        return ranked
